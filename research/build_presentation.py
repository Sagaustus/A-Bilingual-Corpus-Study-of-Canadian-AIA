"""
Build CSDH 2026 conference presentation
20 slides · 20 minutes
Output: assets_for_paper/CSDH2026_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, pathlib

OUT = pathlib.Path(__file__).parent.parent / "assets_for_paper" / "CSDH2026_Presentation.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x3A, 0x5C)   # slide headers / dark bg
CREAM   = RGBColor(0xFB, 0xF8, 0xF2)   # title-slide body
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BURGUN  = RGBColor(0x8B, 0x1A, 0x1A)   # accents / emphasis
TEAL    = RGBColor(0x1A, 0x6B, 0x8A)   # secondary accent
LGRAY   = RGBColor(0xF0, 0xF2, 0xF5)   # content-slide background
MGRAY   = RGBColor(0x9E, 0xA3, 0xAB)   # muted text
DKTEXT  = RGBColor(0x1C, 0x1C, 0x2E)   # body text on light bg
GOLD    = RGBColor(0xC9, 0x9A, 0x06)   # warning / highlight

W  = Inches(13.33)
H  = Inches(7.5)
M  = Inches(0.45)   # margin

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]   # blank layout


# ── Low-level helpers ─────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, l, t, w, h, fill=None, line=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DKTEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return txb

def add_para(tf, text, size=16, bold=False, color=DKTEXT,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    from pptx.util import Pt as PT
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = PT(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = PT(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return p

def header_band(slide, title, subtitle=None):
    """Top navy band with title."""
    add_rect(slide, 0, 0, W, Inches(1.35), fill=NAVY)
    add_text(slide, title,
             M, Inches(0.12), W - 2*M, Inches(0.85),
             size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 M, Inches(0.90), W - 2*M, Inches(0.42),
                 size=14, color=RGBColor(0xB0, 0xC4, 0xDE), italic=True)

def content_box(slide, left=None, top=None, width=None, height=None):
    """Returns a text-frame-ready textbox below the header."""
    l = left  if left   is not None else M
    t = top   if top    is not None else Inches(1.5)
    w = width if width  is not None else W - 2*M
    h = height if height is not None else Inches(5.6)
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.text_frame.word_wrap = True
    return txb.text_frame

def slide_number(slide, n):
    add_text(slide, f"{n} / 20",
             W - Inches(1.1), H - Inches(0.35), Inches(0.9), Inches(0.3),
             size=10, color=MGRAY, align=PP_ALIGN.RIGHT)

def new_slide(n, title, subtitle=None, bg=LGRAY):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, fill=bg)
    header_band(s, title, subtitle)
    slide_number(s, n)
    return s

def bullet(tf, text, level=0, size=17, color=DKTEXT, bold=False,
           marker="▸ ", space_before=4, italic=False):
    indent = "    " * level
    add_para(tf, f"{indent}{marker}{text}",
             size=size, color=color, bold=bold, space_before=space_before,
             italic=italic)

def stat_box(slide, label, value, l, t, w=Inches(2.8), h=Inches(1.4),
             bg=NAVY, val_color=WHITE, lbl_color=RGBColor(0xB0, 0xC4, 0xDE)):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, value, l, t + Inches(0.1), w, Inches(0.75),
             size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t + Inches(0.85), w, Inches(0.5),
             size=12, color=lbl_color, align=PP_ALIGN.CENTER)

def divider(slide, t=None, color=BURGUN):
    top = t if t else Inches(1.38)
    add_rect(slide, M, top, Inches(0.06), Inches(4.5), fill=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, W, H, fill=NAVY)
# decorative stripe
add_rect(s1, 0, Inches(5.1), W, Inches(0.08), fill=BURGUN)
add_rect(s1, 0, Inches(5.18), W, Inches(2.32), fill=RGBColor(0x12, 0x28, 0x42))

add_text(s1, "CSDH / SCHN 2026  ·  Untranslatable",
         M, Inches(0.25), W - 2*M, Inches(0.4),
         size=13, color=RGBColor(0x7A, 0xA8, 0xCC), italic=True)

add_text(s1,
         "Abstract Rhetorics of Algorithmic Governance",
         M, Inches(0.75), W - 2*M, Inches(1.5),
         size=38, bold=True, color=WHITE)

add_text(s1,
         "A Bilingual Corpus Study of Canadian Algorithmic Impact Assessments",
         M, Inches(2.2), W - 2*M, Inches(0.9),
         size=22, color=RGBColor(0xB0, 0xC4, 0xDE), italic=True)

add_text(s1,
         "47 AIAs  ·  114 source files  ·  1,178 structured records  ·  37 ethical terms  ·  1,130 corpus occurrences",
         M, Inches(3.3), W - 2*M, Inches(0.5),
         size=13, color=RGBColor(0x7A, 0xA8, 0xCC))

add_text(s1, "[Author Name]  |  [Institution]",
         M, Inches(5.4), Inches(6), Inches(0.4),
         size=16, color=WHITE, bold=True)
add_text(s1, "20-minute paper  ·  CSDH/SCHN 2026",
         M, Inches(5.82), Inches(6), Inches(0.4),
         size=13, color=RGBColor(0x7A, 0xA8, 0xCC))

slide_number(s1, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Question
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide(2, "The Question", "What bilingualism means for algorithmic governance")

add_text(s2,
    "“Can Francophone Canadians read, in their language, the government’s "
    "own account of how an algorithm might affect their welfare benefit, "
    "their immigration application, their disability claim?”",
    Inches(1.2), Inches(1.55), Inches(11.0), Inches(2.1),
    size=22, italic=True, color=NAVY, bold=False)

tf = content_box(s2, top=Inches(3.85), height=Inches(3.2))
bullet(tf, "Directive on Automated Decision-Making (TBS, 2019) — AIAs mandatory in both official languages", size=16)
bullet(tf, "Official Languages Act (1985) — constitutional entitlement, not a courtesy", size=16)
bullet(tf, "35% of Canadians are primarily Francophone; many exercise rights in French in immigration, welfare, disability", size=16)
bullet(tf, "This paper asks: does the bilingual guarantee hold when algorithms govern?", size=16, bold=True, color=BURGUN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is an AIA?
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide(3, "What is an Algorithmic Impact Assessment?")

tf = content_box(s3, top=Inches(1.55), height=Inches(2.6))
bullet(tf, "Structured self-assessment required before deploying AI in public decision-making", size=16)
bullet(tf, "Covers: system purpose · data inputs · individual rights impacts · safeguards · justification for automation", size=16)
bullet(tf, "Numerical scores → Impact Level I–IV (higher = more scrutiny, external review, public disclosure)", size=16)
bullet(tf, "Narrative free-text fields provide the substance of accountability that scores alone cannot convey", size=16, bold=True)

# Key corpus facts
stat_box(s3, "AIAs harvested", "47", Inches(0.45), Inches(4.2), bg=NAVY)
stat_box(s3, "Structured submissions", "30", Inches(3.45), Inches(4.2), bg=TEAL)
stat_box(s3, "Bilingual PDF pairs", "16", Inches(6.45), Inches(4.2), bg=BURGUN)
stat_box(s3, "Source files", "114", Inches(9.45), Inches(4.2), bg=RGBColor(0x3A, 0x5A, 0x3A))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Corpus & Methodology
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide(4, "Corpus & Methodology", "Three analytical layers · Two LLMs · One database")

# Left column
add_text(s4, "Three Corpus Layers", M, Inches(1.55), Inches(6), Inches(0.4),
         size=15, bold=True, color=NAVY)
tf_l = content_box(s4, left=M, top=Inches(2.0), width=Inches(5.9), height=Inches(3.5))
bullet(tf_l, "Layer 1 — Full 47 AIAs: discoverability & availability audit", size=14)
bullet(tf_l, "Layer 2 — 30 structured submissions: PostgreSQL database (27 tables, 1,178 rows), LLM semantic interpretation", size=14)
bullet(tf_l, "Layer 3 — 16 bilingual PDF pairs: governance terminology extraction (21 terms, 751 occurrences)", size=14)

# Right column
add_text(s4, "Two-Model LLM Pipeline", Inches(7.0), Inches(1.55), Inches(5.9), Inches(0.4),
         size=15, bold=True, color=NAVY)
tf_r = content_box(s4, left=Inches(7.0), top=Inches(2.0), width=Inches(5.9), height=Inches(3.5))
bullet(tf_r, "Llama 3.3 70B (IONOS cloud) — Layer 2: bilingual divergence, automation justification, risk/rights, safeguard compliance", size=14)
bullet(tf_r, "GPT-4 Turbo — Layer 3: governance terminology extraction + semantic grouping across 32 documents", size=14)
bullet(tf_r, "Validation: LLM risk labels vs. computed scores (near-perfect ordinal consistency)", size=14)
bullet(tf_r, "Inter-rater: κ = 0.698 (substantial agreement, Landis & Koch 1977)", size=14)

add_rect(s4, Inches(6.75), Inches(1.55), Inches(0.04), Inches(5.6), fill=MGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Three Failure Modes (overview)
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide(5, "Three Failure Modes", "The answer to our question — in three acts")

for i, (icon, label, desc, col) in enumerate([
    ("①", "Discoverability Gap",
     "French versions exist but are structured for inaccessibility — only 4% detectable by metadata search",
     NAVY),
    ("②", "Availability Gap",
     "23% of systems are genuinely monolingual — no French documentation exists in any form",
     BURGUN),
    ("③", "Semantic Drift",
     "96.7% of bilingual pairs diverge — primarily through omission, then conceptual fracture",
     TEAL),
]):
    l = M + i * Inches(4.35)
    add_rect(s5, l, Inches(1.55), Inches(4.1), Inches(5.1), fill=col)
    add_text(s5, icon, l, Inches(1.65), Inches(4.1), Inches(0.9),
             size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s5, label, l, Inches(2.55), Inches(4.1), Inches(0.55),
             size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s5, l + Inches(0.2), Inches(3.15), Inches(3.7), Inches(0.04),
             fill=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s5, desc, l + Inches(0.15), Inches(3.25), Inches(3.8), Inches(3.1),
             size=14, color=CREAM, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Discoverability Gap
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide(6, "① Discoverability Gap", "The invisible French archive")

add_text(s6, "4%", Inches(0.45), Inches(1.6), Inches(2.8), Inches(1.4),
         size=72, bold=True, color=BURGUN, align=PP_ALIGN.CENTER)
add_text(s6, "initially detectable\nas having FR counterparts",
         Inches(0.45), Inches(3.0), Inches(2.8), Inches(0.8),
         size=14, color=DKTEXT, align=PP_ALIGN.CENTER)

add_rect(s6, Inches(3.5), Inches(2.2), Inches(0.6), Inches(0.6), fill=BURGUN)
add_text(s6, "→", Inches(3.5), Inches(2.1), Inches(0.6), Inches(0.6),
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s6, "38%", Inches(4.3), Inches(1.6), Inches(2.8), Inches(1.4),
         size=72, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s6, "after metadata\ncorrection (9.5× increase)",
         Inches(4.3), Inches(3.0), Inches(2.8), Inches(0.8),
         size=14, color=DKTEXT, align=PP_ALIGN.CENTER)

tf6 = content_box(s6, left=Inches(7.5), top=Inches(1.6), width=Inches(5.4), height=Inches(4.0))
bullet(tf6, "EN titles: [System Name] — [Department]  (standard form)", size=14)
bullet(tf6, "FR titles: abbreviated, grammatically inverted, or absent from dataset identifier", size=14)
bullet(tf6, "A Francophone searching open.canada.ca in standard terms would find fewer than 1 in 20", size=14, bold=True, color=BURGUN)
add_para(tf6, "", size=8)
bullet(tf6, "Discoverability failure is the first mechanism through which the law's requirements are structurally circumvented", size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Availability Gap
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide(7, "② Availability Gap", "Structural monolingualism")

stat_box(s7, "systems with NO French documentation", "23%",
         Inches(0.6), Inches(1.7), w=Inches(3.5), h=Inches(1.8), bg=BURGUN)

tf7 = content_box(s7, left=Inches(4.5), top=Inches(1.65), width=Inches(8.4), height=Inches(2.0))
bullet(tf7, "Not inadequately translated — French documentation was never produced", size=16, bold=True)
bullet(tf7, "Highest rates: CBSA and IRCC — agencies whose systems most directly affect Francophone applicants", size=16)
bullet(tf7, "Structural monolingualism is highest where bilingual accountability is most needed", size=16, color=BURGUN, bold=True)

# Bottom callout
add_rect(s7, M, Inches(4.0), W - 2*M, Inches(2.9), fill=NAVY)
add_text(s7,
    "The riskiest systems have the thinnest governance documentation in every dimension simultaneously: "
    "maximum risk scores  ·  NULL automation type disclosure  ·  no trade-off analysis  ·  no French content",
    Inches(0.7), Inches(4.2), W - Inches(1.4), Inches(2.6),
    size=16, color=WHITE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Semantic Drift: Omission
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide(8, "③ Semantic Drift — The Architecture of Omission")

for val, lbl, bg, l in [
    ("96.7%", "of bilingual pairs diverge", NAVY, Inches(0.4)),
    ("2.61 / 5", "average fidelity score", TEAL, Inches(3.55)),
    ("71%", "of divergence = omission", BURGUN, Inches(6.7)),
    ("43%", "score 0 or 1 (effectively monolingual)", RGBColor(0x4A, 0x4A, 0x7A), Inches(9.85)),
]:
    add_rect(s8, l, Inches(1.6), Inches(3.0), Inches(1.5), fill=bg)
    add_text(s8, val, l, Inches(1.62), Inches(3.0), Inches(0.9),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s8, lbl, l, Inches(2.5), Inches(3.0), Inches(0.55),
             size=11, color=CREAM, align=PP_ALIGN.CENTER)

tf8 = content_box(s8, top=Inches(3.3), height=Inches(3.7))
bullet(tf8, "The fidelity distribution is binary, not gradual — an institutional cliff: departments either treat French as essential or they do not", size=15, bold=True)
bullet(tf8, "Only one submission (family reunification, submission 112) achieves full bilingual equivalence — score 5/5", size=15)
bullet(tf8, "Reframing: 9%  ·  Translation differences: 9%  ·  Additions in French: 8%  ·  Terminological: 3%", size=15)
bullet(tf8, "The French version of Canadian algorithmic governance is not poorly translated — it is largely absent", size=15, italic=True, color=BURGUN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — What Gets Omitted (table)
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide(9, "What Gets Omitted — Top Divergent Fields")

# Table
rows, cols = 8, 3
left, top, width, height = Inches(0.45), Inches(1.6), Inches(12.4), Inches(4.8)
tbl = s9.shapes.add_table(rows, cols, left, top, width, height).table

headers = ["Field", "Omissions", "Governance Function"]
col_widths = [Inches(4.8), Inches(1.6), Inches(6.0)]
for i, cw in enumerate(col_widths):
    tbl.columns[i].width = cw

data = [
    ("Evaluation criteria — how the system judges people",     "14", "Most accountability-critical"),
    ("System output description — what the system produces",   "12", "Foundation of public understanding"),
    ("Project title — what the system is called",              "12", "Discoverability prerequisite"),
    ("System description — what the system does",              "12", "Basic transparency requirement"),
    ("Rights and freedoms assessment",                         "9",  "Charter obligations"),
    ("Equality and dignity assessment",                        "9",  "Human rights framework"),
    ("Expected improvements — claimed public benefits",        "6",  "Justification for automation"),
]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.runs[0].font.color.rgb = WHITE
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(13)
    p.runs[0].font.name = "Calibri"

for ri, row_data in enumerate(data):
    for ci, val in enumerate(row_data):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LGRAY
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.name = "Calibri"
        p.runs[0].font.color.rgb = DKTEXT
        if ci == 1:
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = BURGUN
            p.alignment = PP_ALIGN.CENTER

add_text(s9, "Every top-10 divergent field is an omission. Source: interp_bilingual_divergence, 30 submissions.",
         M, Inches(6.5), W - 2*M, Inches(0.35),
         size=11, color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — The Control Group: What Translates Cleanly
# ══════════════════════════════════════════════════════════════════════════════
s10 = new_slide(10, "The Control Group — What Translates Cleanly",
                "Abstract ethical principles translate at near-perfect parity")

# Four boxes
for val, en, fr, col, l in [
    ("≈ 1.00", "fairness",       "équité",        TEAL,   Inches(0.45)),
    ("1.00",   "transparency",   "transparence",  NAVY,   Inches(3.75)),
    ("1.00",   "discrimination", "discrimination",TEAL,   Inches(7.05)),
    ("1.00",   "explanation",    "explication",   NAVY,   Inches(10.35)),
]:
    add_rect(s10, l, Inches(1.65), Inches(2.8), Inches(2.2), fill=col)
    add_text(s10, val, l, Inches(1.7), Inches(2.8), Inches(0.9),
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s10, f"{en} / {fr}", l, Inches(2.6), Inches(2.8), Inches(0.65),
             size=13, color=CREAM, align=PP_ALIGN.CENTER, italic=True)
    add_text(s10, "EN÷FR ratio", l, Inches(3.25), Inches(2.8), Inches(0.4),
             size=10, color=RGBColor(0xA0, 0xC0, 0xD0), align=PP_ALIGN.CENTER)

tf10 = content_box(s10, top=Inches(4.1), height=Inches(3.0))
bullet(tf10, "Parity ≈ 1.0 proves divergence in other terms is not an artifact of French being a different language", size=16, bold=True)
bullet(tf10, "Not word-count differences. Not random noise. The divergence in accountability, oversight, audit, review is specific, patterned, structural", size=16)
bullet(tf10, "These shared concepts come from traditions common to both Common Law and Civil Law — the untranslatable is what doesn't cross that boundary", size=16)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — The Accountability Fracture
# ══════════════════════════════════════════════════════════════════════════════
s11 = new_slide(11, "The Untranslatable I — The Accountability Fracture")

add_text(s11, '"accountability"', M, Inches(1.6), Inches(4.5), Inches(0.7),
         size=30, bold=True, color=NAVY, italic=True)
add_text(s11, "15× in English AIA documents", M, Inches(2.3), Inches(4.5), Inches(0.5),
         size=15, color=MGRAY)
add_text(s11, "0× direct equivalent in French", M, Inches(2.8), Inches(4.5), Inches(0.5),
         size=15, color=BURGUN, bold=True)

# Arrow
add_text(s11, "→  fractures into:", Inches(4.3), Inches(2.3), Inches(2.0), Inches(0.5),
         size=17, color=NAVY, bold=True)

# Three boxes
for label, freq, desc, col, t in [
    ("reddition de comptes", "11×", "Administrative reporting\n(Federal public administration)", NAVY,    Inches(1.7)),
    ("imputabilité",         "4×",  "Technical attribution\n(Civil law, causal liability)",     TEAL,    Inches(3.05)),
    ("responsabilité",       "1×",  "Legal/moral responsibility\n(Both traditions)",            BURGUN,  Inches(4.4)),
]:
    add_rect(s11, Inches(6.5), t, Inches(6.38), Inches(1.2), fill=col)
    add_text(s11, f"« {label} »  {freq}", Inches(6.65), t + Inches(0.08), Inches(6.0), Inches(0.45),
             size=16, bold=True, color=WHITE)
    add_text(s11, desc, Inches(6.65), t + Inches(0.52), Inches(6.0), Inches(0.6),
             size=12, color=CREAM, italic=True)

add_text(s11,
    "English uses a single superordinate spanning legal, administrative, and technical domains.\n"
    "French distributes it across terms with different institutional weights — a governance decision, "
    "not a translation error.\nThis is Cassin's untranslatable: a term one does not cease (not) to translate.",
    M, Inches(5.7), W - 2*M, Inches(1.5),
    size=13, color=DKTEXT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Oversight & Bias
# ══════════════════════════════════════════════════════════════════════════════
s12 = new_slide(12, "The Untranslatable II — Oversight & Bias")

# Oversight section
add_text(s12, "OVERSIGHT", M, Inches(1.6), Inches(6.2), Inches(0.45),
         size=14, bold=True, color=BURGUN)
add_rect(s12, M, Inches(2.05), Inches(5.9), Inches(1.7), fill=NAVY)
add_text(s12, '"monitoring"   5×  EN', Inches(0.6), Inches(2.1), Inches(5.3), Inches(0.55),
         size=18, bold=True, color=WHITE)
add_text(s12, 'collocates: regular · quality · assurance · measures → observational, passive',
         Inches(0.6), Inches(2.65), Inches(5.3), Inches(0.55),
         size=12, color=CREAM, italic=True)

add_text(s12, '"contrôle"   27×  FR', Inches(0.6), Inches(3.82), Inches(5.3), Inches(0.55),
         size=18, bold=True, color=WHITE)
add_rect(s12, M, Inches(3.75), Inches(5.9), Inches(1.7), fill=TEAL)
add_text(s12, 'collocates: gouvernement · données · fédéral · points → regulatory control, active authority',
         Inches(0.6), Inches(4.35), Inches(5.3), Inches(0.55),
         size=12, color=CREAM, italic=True)

add_text(s12, "5.4× ratio — the largest divergence in the corpus",
         M, Inches(5.52), Inches(6.0), Inches(0.4),
         size=13, color=BURGUN, bold=True)

add_rect(s12, Inches(6.35), Inches(1.55), Inches(0.04), Inches(5.1), fill=MGRAY)

# Bias section
add_text(s12, "BIAS", Inches(6.6), Inches(1.6), Inches(6.2), Inches(0.45),
         size=14, bold=True, color=BURGUN)

add_rect(s12, Inches(6.6), Inches(2.05), Inches(6.3), Inches(0.85), fill=NAVY)
add_text(s12, '"bias"  34×  EN  — one term, two registers collapsed',
         Inches(6.8), Inches(2.1), Inches(6.0), Inches(0.7),
         size=15, bold=True, color=WHITE)

add_rect(s12, Inches(6.6), Inches(2.95), Inches(3.0), Inches(1.5), fill=TEAL)
add_text(s12, "biais   22×", Inches(6.75), Inches(3.0), Inches(2.7), Inches(0.5),
         size=16, bold=True, color=WHITE)
add_text(s12, "Technical/algorithmic\n→ unintended model output", Inches(6.75), Inches(3.5), Inches(2.7), Inches(0.85),
         size=12, color=CREAM, italic=True)

add_rect(s12, Inches(9.8), Inches(2.95), Inches(3.1), Inches(1.5), fill=BURGUN)
add_text(s12, "préjugé   9×", Inches(9.95), Inches(3.0), Inches(2.8), Inches(0.5),
         size=16, bold=True, color=WHITE)
add_text(s12, "Social/moral prejudice\n→ collocates with atténuation", Inches(9.95), Inches(3.5), Inches(2.8), Inches(0.85),
         size=12, color=CREAM, italic=True)

add_text(s12,
    "French makes a lexical distinction English collapses — a conceptual refinement, not a translation error.\n"
    "Francophone authors distinguish more precisely between a data bias in the model and a human prejudice in the process.",
    Inches(6.6), Inches(4.6), Inches(6.3), Inches(0.9),
    size=12, color=DKTEXT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Philosophical Drift (37 terms)
# ══════════════════════════════════════════════════════════════════════════════
s13 = new_slide(13, "Philosophical Drift — Systematic Evidence Across 37 Terms",
                "Applied Ethics · Normative Ethics · Metaethics  ·  1,130 corpus occurrences")

# EN bar
add_rect(s13, M, Inches(1.65), Inches(4.6), Inches(0.55), fill=BURGUN)
add_text(s13, "EN  INSTRUMENTALIZED  37.8%  (14/37 terms)",
         Inches(0.6), Inches(1.68), Inches(9), Inches(0.5),
         size=14, bold=True, color=WHITE)

add_rect(s13, M, Inches(2.3), Inches(2.3), Inches(0.55), fill=RGBColor(0x9B, 0x59, 0x59))
add_text(s13, "EN  narrowed  18.9%",
         Inches(0.6), Inches(2.33), Inches(9), Inches(0.5),
         size=13, color=WHITE)

add_rect(s13, M, Inches(2.95), Inches(1.65), Inches(0.55), fill=RGBColor(0x7A, 0x7A, 0x9A))
add_text(s13, "EN  reframed  13.5%",
         Inches(0.6), Inches(2.98), Inches(9), Inches(0.5),
         size=13, color=WHITE)

# FR bar
add_rect(s13, M, Inches(3.7), Inches(3.6), Inches(0.55), fill=TEAL)
add_text(s13, "FR  REFRAMED  29.7%  (11/37 terms)",
         Inches(0.6), Inches(3.73), Inches(9), Inches(0.5),
         size=14, bold=True, color=WHITE)

add_rect(s13, M, Inches(4.35), Inches(2.0), Inches(0.55), fill=RGBColor(0x8B, 0x3A, 0x3A))
add_text(s13, "FR  absent  21.6%",
         Inches(0.6), Inches(4.38), Inches(9), Inches(0.5),
         size=13, color=WHITE)

add_rect(s13, M, Inches(5.0), Inches(1.9), Inches(0.55), fill=RGBColor(0x5A, 0x8A, 0x6A))
add_text(s13, "FR  instrumentalized  16.2%",
         Inches(0.6), Inches(5.03), Inches(9), Inches(0.5),
         size=13, color=WHITE)

# Stats panel
add_rect(s13, Inches(7.9), Inches(1.65), Inches(5.0), Inches(4.7), fill=NAVY)
tf13 = content_box(s13, left=Inches(8.1), top=Inches(1.75), width=Inches(4.6), height=Inches(4.5))
add_para(tf13, "Statistical Evidence", size=15, bold=True, color=WHITE)
add_para(tf13, "", size=6)
bullet(tf13, "Fisher's exact: EN instrumentalizes > FR", size=13, color=CREAM, marker="")
bullet(tf13, "OR = 2.99,  p = .033", size=15, bold=True, color=GOLD, marker="  ")
add_para(tf13, "", size=4)
bullet(tf13, "FR reframes > EN (trend)", size=13, color=CREAM, marker="")
bullet(tf13, "OR = 2.56,  p = .078", size=14, color=CREAM, marker="  ")
add_para(tf13, "", size=4)
bullet(tf13, "Cramér's V = 0.322 (moderate effect)", size=13, color=CREAM, marker="")
add_para(tf13, "", size=4)
bullet(tf13, "Total EN (665) vs FR (465)", size=13, color=CREAM, marker="")
bullet(tf13, "Wilcoxon W = 496,  p = .001", size=14, bold=True, color=GOLD, marker="  ")
add_para(tf13, "", size=4)
bullet(tf13, "Spearman r = 0.817,  p < .001", size=13, color=CREAM, marker="")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Three Conceptual Shifts
# ══════════════════════════════════════════════════════════════════════════════
s14 = new_slide(14, "The Deepest Divergence — Three Conceptual Shifts",
                "Terms that encode categorically different conceptual orientations")

add_text(s14,
    "For these three terms, EN and FR do not differ in degree — they differ in kind.",
    M, Inches(1.6), W - 2*M, Inches(0.45),
    size=16, italic=True, color=NAVY)

for i, (term, en_drift, fr_drift, contrast) in enumerate([
    ("accountability",
     "narrowed  (individual assignment, Common Law)",
     "absent  (distributed: reddition / imputabilité / responsabilité)",
     "EN constructs accountability as a named personal property. FR has no single equivalent — the concept is juridically untranslatable."),
    ("oversight",
     "instrumentalized  (monitoring = passive quality assurance)",
     "reframed  (contrôle = active regulatory authority)",
     "Same AIA question. Two different theories of the state's role in algorithmic governance."),
    ("justice",
     "narrowed  (procedural fairness only)",
     "expanded  (social justice, distributive obligation — Civil Law tradition)",
     "The only term where French is philosophically richer than English."),
]):
    t = Inches(2.15) + i * Inches(1.68)
    add_rect(s14, M, t, W - 2*M, Inches(1.58), fill=LGRAY)
    add_rect(s14, M, t, Inches(0.07), Inches(1.58), fill=BURGUN)
    add_text(s14, f"«  {term}  »", Inches(0.65), t + Inches(0.06), Inches(3.0), Inches(0.5),
             size=17, bold=True, color=NAVY)
    add_text(s14, f"EN: {en_drift}", Inches(3.85), t + Inches(0.06), Inches(5.1), Inches(0.42),
             size=13, color=BURGUN, bold=True)
    add_text(s14, f"FR: {fr_drift}", Inches(3.85), t + Inches(0.5), Inches(8.8), Inches(0.42),
             size=13, color=TEAL, bold=True)
    add_text(s14, contrast, Inches(0.65), t + Inches(0.96), W - Inches(1.1), Inches(0.55),
             size=12, color=DKTEXT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Where FR is More Faithful
# ══════════════════════════════════════════════════════════════════════════════
s15 = new_slide(15, "The Counterintuitive Finding — Where French Is More Faithful")

add_text(s15,
    "Two terms where French preserves the philosophical register that English abandons:",
    M, Inches(1.6), W - 2*M, Inches(0.45),
    size=16, bold=True, color=NAVY)

for l, term, en_d, fr_d, detail, col in [
    (Inches(0.45), "recourse", "instrumentalized", "✅ faithful",
     "English reduces recourse to procedural mechanism — appeals forms, complaint channels.\n"
     "French preserves recourse as a citizen's normative right to remedy.",
     TEAL),
    (Inches(6.85), "non-maleficence", "👻 unnamed", "✅ faithful",
     "English suppresses the philosophical label — harm prevention is present but unnamed.\n"
     "French frames ne pas nuire as an ethical principle, not just a risk-management task.",
     NAVY),
]:
    add_rect(s15, l, Inches(2.1), Inches(6.2), Inches(3.85), fill=col)
    add_text(s15, f"« {term} »", l + Inches(0.2), Inches(2.2), Inches(5.8), Inches(0.65),
             size=22, bold=True, color=WHITE)
    add_text(s15, f"EN: {en_d}", l + Inches(0.2), Inches(2.9), Inches(5.8), Inches(0.45),
             size=14, color=CREAM)
    add_text(s15, f"FR: {fr_d}", l + Inches(0.2), Inches(3.35), Inches(5.8), Inches(0.45),
             size=14, color=GOLD, bold=True)
    add_text(s15, detail, l + Inches(0.2), Inches(3.82), Inches(5.8), Inches(1.9),
             size=12, color=CREAM, italic=True)

add_rect(s15, M, Inches(6.1), W - 2*M, Inches(0.92), fill=BURGUN)
add_text(s15,
    "Both terms concern citizen remedy and harm prevention — the dimensions most directly relevant to "
    "Francophone Canadians' rights. Their preservation in French, against the general pattern, is the "
    "corpus's most counterintuitive finding.",
    Inches(0.65), Inches(6.15), W - Inches(1.3), Inches(0.85),
    size=13, color=WHITE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Automation Rhetoric
# ══════════════════════════════════════════════════════════════════════════════
s16 = new_slide(16, "Automation Rhetoric — The Universal Pattern",
                "Across all 30 systems, in all 6 sectors: automation is invariably framed as 'assistance'")

# Quote strip
add_rect(s16, 0, Inches(1.6), W, Inches(1.1), fill=NAVY)
add_text(s16,
    '"The system  helps sort and assign  applications...  all decisions are  made by an officer."  '
    '|  "a tool  for the purpose of assisting..."  decisions may be rendered  without direct human involvement.',
    Inches(0.5), Inches(1.65), W - Inches(1.0), Inches(0.95),
    size=13, italic=True, color=CREAM)

tf16 = content_box(s16, top=Inches(2.85), height=Inches(2.4))
add_para(tf16, "Four rhetorical strategies operating across the corpus:", size=15, bold=True, color=NAVY)
add_para(tf16, "", size=5)
bullet(tf16, "Logistical sanitization — automation as operational necessity ('managing growing volumes'), preventing scrutiny", size=14)
bullet(tf16, "Clerical reductionism — consequential decisions described as 'administrative groupings' or 'sorting mechanisms'", size=14)
bullet(tf16, "Technological diminishment — algorithm's power minimized ('assists,' 'supports,' 'recommends')", size=14)
bullet(tf16, "Procedural insulation — safeguard claims answered with bare 'Yes' without elaboration", size=14)

# Contradiction callout
add_rect(s16, M, Inches(5.35), W - 2*M, Inches(1.7), fill=BURGUN)
add_text(s16,
    "Passport Facial Recognition System — same document, same paragraph:\n"
    '"decisions may be rendered without direct human involvement"  &  '
    '"the human-in-the-loop remains the definitive source of judgment and accountability"\n'
    "The AIA does not resolve this contradiction — it contains it.",
    Inches(0.65), Inches(5.4), W - Inches(1.3), Inches(1.6),
    size=13, color=WHITE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Safeguard Compliance Funnel
# ══════════════════════════════════════════════════════════════════════════════
s17 = new_slide(17, "The Performative Compliance Funnel",
                "From universal claim to near-total non-disclosure")

# Funnel visual
for pct, label, col, t, w in [
    ("100%", "Human override enabled (claimed)",     TEAL,                 Inches(1.65), Inches(12.4)),
    ("67%",  "GBA+ conducted",                       RGBColor(0x2A,0x7A,0x6A), Inches(2.45), Inches(10.0)),
    ("60%",  "Bias testing documented",              RGBColor(0xC0,0x70,0x30), Inches(3.25), Inches(7.5)),
    ("3%",   "Bias testing results public",          BURGUN,               Inches(4.05), Inches(3.0)),
]:
    lft = (W - w) / 2
    add_rect(s17, lft, t, w, Inches(0.68), fill=col)
    add_text(s17, f"{pct}  —  {label}",
             lft + Inches(0.2), t + Inches(0.08), w - Inches(0.4), Inches(0.55),
             size=14, bold=(pct in ["100%", "3%"]), color=WHITE)

add_text(s17,
    "The form is completed. The box is checked. The record shows compliance.\n"
    "The downstream requirements that would make compliance meaningful are not fulfilled.",
    M, Inches(5.0), W - 2*M, Inches(0.8),
    size=15, italic=True, color=NAVY, bold=False)

add_rect(s17, M, Inches(5.9), W - 2*M, Inches(1.15), fill=NAVY)
add_text(s17,
    "The AIA asks 'have you done X?' rather than 'show how you did X.'\n"
    "The instrument elicits claims of compliance rather than evidence of it.",
    Inches(0.65), Inches(5.95), W - Inches(1.3), Inches(1.05),
    size=14, color=WHITE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Translation as Governance Transformation
# ══════════════════════════════════════════════════════════════════════════════
s18 = new_slide(18, "Translation as Governance Transformation",
                "Baker's non-equivalence thesis — structural, not merely lexical")

tf18_l = content_box(s18, left=M, top=Inches(1.6), width=Inches(6.1), height=Inches(4.6))
add_para(tf18_l, "The Clean / Untranslatable Pattern", size=15, bold=True, color=NAVY)
add_para(tf18_l, "", size=5)
bullet(tf18_l, "Abstract principles (fairness, transparency, discrimination) — shared across legal traditions — translate at ratio ≈1.0", size=14)
bullet(tf18_l, "Institutional practice terms (accountability, oversight, audit) — encoding Common Law assumptions — diverge dramatically", size=14)
bullet(tf18_l, "Bilingual failure is not a function of linguistic difference — it is a function of institutional culture", size=14, bold=True, color=BURGUN)
add_para(tf18_l, "", size=5)
add_para(tf18_l, "The Francophone Accountability Deficit", size=15, bold=True, color=NAVY)
add_para(tf18_l, "", size=5)
bullet(tf18_l, "43% of submissions: structurally monolingual — French treated as non-essential", size=14)
bullet(tf18_l, "57% genuinely bilingual — no gradation, an institutional cliff", size=14)
bullet(tf18_l, "Venuti's domestication: dominant language asserts itself through absence of translation", size=14, italic=True)

add_rect(s18, Inches(6.35), Inches(1.55), Inches(0.04), Inches(5.3), fill=MGRAY)

tf18_r = content_box(s18, left=Inches(6.6), top=Inches(1.6), width=Inches(6.3), height=Inches(5.3))
add_para(tf18_r, "Cassin's Untranslatable — Measured", size=15, bold=True, color=NAVY)
add_para(tf18_r, "", size=5)
bullet(tf18_r, '"Accountability" is not untranslatable because it cannot be rendered in French', size=14)
bullet(tf18_r, "It is untranslatable because no single French term covers the same conceptual territory — and the translation choice is itself a governance act", size=14, bold=True)
add_para(tf18_r, "", size=8)
add_para(tf18_r,
    '"The English AIA can ask \'have you assigned accountability?\' as a unified governance question. The French AIA cannot — it must choose which kind of accountability it is asking about."',
    size=14, italic=True, color=TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Policy + DH Contributions
# ══════════════════════════════════════════════════════════════════════════════
s19 = new_slide(19, "Implications — Policy & Digital Humanities")

# Left: Policy
add_text(s19, "Policy Recommendations", M, Inches(1.6), Inches(6.1), Inches(0.4),
         size=15, bold=True, color=NAVY)
tf19_l = content_box(s19, left=M, top=Inches(2.05), width=Inches(6.1), height=Inches(4.5))
bullet(tf19_l, "① Bilingual field completion as a condition of AIA publication — the pipeline should reject submissions with NULL French narrative fields in rights, description, evaluation criteria", size=13)
bullet(tf19_l, "② Fix the Impact Level mechanism — all 30 interpreted submissions have NULL impact_level, so the escalating scrutiny requirements were never triggered", size=13)
bullet(tf19_l, "③ Mandatory publication of GBA+ and bias testing results — not just their completion", size=13)
add_para(tf19_l, "", size=6)
bullet(tf19_l, "Submission 112 proves bilingual AIA governance is achievable — it is an institutional choice, not a linguistic impossibility", size=13, italic=True, color=BURGUN)

add_rect(s19, Inches(6.35), Inches(1.55), Inches(0.04), Inches(5.3), fill=MGRAY)

# Right: DH
add_text(s19, "DH Contributions", Inches(6.6), Inches(1.6), Inches(6.3), Inches(0.4),
         size=15, bold=True, color=NAVY)
tf19_r = content_box(s19, left=Inches(6.6), top=Inches(2.05), width=Inches(6.3), height=Inches(4.5))
bullet(tf19_r, "LLM-assisted semantic auditing as a DH method — structured comparative analysis of bilingual text pairs at corpus scale, where qualitative reading cannot", size=13)
bullet(tf19_r, "Relational corpus design as reproducibility infrastructure — every finding is a SQL query, every LLM output is structured JSON, every claim is independently verifiable", size=13)
bullet(tf19_r, "AI governance as a DH problem — not background context but the central site where computational text analysis, corpus linguistics, and the cultural analysis of institutional documents meet questions of power and justice", size=13, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s20 = prs.slides.add_slide(BLANK)
add_rect(s20, 0, 0, W, H, fill=NAVY)
add_rect(s20, 0, Inches(4.9), W, Inches(0.08), fill=BURGUN)
add_rect(s20, 0, Inches(4.98), W, Inches(2.52), fill=RGBColor(0x12, 0x28, 0x42))

add_text(s20, "Conclusion", M, Inches(0.3), W - 2*M, Inches(0.55),
         size=14, color=RGBColor(0x7A, 0xA8, 0xCC), bold=True)

add_text(s20,
    "Canada's bilingual AIA corpus reveals a state that is not untranslatable —\n"
    "submission 112 proves bilingual algorithmic governance is achievable —\n"
    "but a state that does not translate, systematically, the governance content that matters most.",
    M, Inches(0.95), W - 2*M, Inches(1.55),
    size=20, color=WHITE, bold=False, italic=True)

add_rect(s20, M, Inches(2.6), W - 2*M, Inches(0.04), fill=RGBColor(0x7A, 0xA8, 0xCC))

# Three takeaways
for i, txt in enumerate([
    "The untranslatable state is primarily the untranslated state",
    "English instrumentalizes; French reframes — but both do so within a framework of assumed values rather than reasoned ones",
    "Quantifying the francophone accountability deficit is not a neutral operation — it is a precondition for accountability",
]):
    add_rect(s20, M, Inches(2.75) + i * Inches(0.68), Inches(0.07), Inches(0.52),
             fill=BURGUN)
    add_text(s20, txt, Inches(0.7), Inches(2.77) + i * Inches(0.68),
             W - Inches(1.1), Inches(0.52),
             size=15, color=CREAM)

add_text(s20,
    '"When algorithms make decisions about refugee applications, disability benefits, and mental health eligibility,\n'
    'the analysis of that text is not peripheral to questions of power and justice — it is central to them."',
    M, Inches(5.05), W - 2*M, Inches(1.1),
    size=14, italic=True, color=RGBColor(0xB0, 0xC4, 0xDE))

add_text(s20,
    "Data & code: github.com/Sagaustus/A-Bilingual-Corpus-Study-of-Canadian-AIA",
    M, Inches(6.25), W - 2*M, Inches(0.4),
    size=12, color=RGBColor(0x7A, 0xA8, 0xCC))

slide_number(s20, 20)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✓ Saved  {OUT}")
print(f"  {len(prs.slides)} slides")
